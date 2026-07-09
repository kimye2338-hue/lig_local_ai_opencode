# -*- coding: utf-8 -*-
"""One-shot pending validation report for OpenCodeLIG user packages.

The target PC is offline and may have company-only applications installed.
This checker is deliberately broad: one BAT should collect enough evidence to
finish the remaining integration work without asking the user to run another
ad-hoc probe.

Safety rules:
- never delete or overwrite OpenCodeLIG_USERDATA;
- never print gateway URL or API key;
- keep live app probes short and non-mutating;
- write one human-readable Markdown report as the primary output.
"""
from __future__ import annotations

import argparse
import hashlib
import importlib.util
import json
import os
import platform
import re
import shutil
import subprocess
import sys
import tempfile
import urllib.request
from dataclasses import asdict, dataclass
from datetime import datetime
from pathlib import Path
from typing import Any, Iterable

_WS_ROOT = Path(__file__).resolve().parents[1]
if str(_WS_ROOT) not in sys.path:
    sys.path.insert(0, str(_WS_ROOT))


@dataclass
class Check:
    section: str
    item: str
    status: str
    evidence: str
    next_action: str = ""


STATUSES = ("PASS", "WARN", "PENDING", "FAIL", "SKIP")
OPTIONAL_MODULE_WHEELS = {
    "openpyxl": ("openpyxl",),
    "docx": ("python_docx", "docx"),
    "pptx": ("python_pptx", "pptx"),
    "markitdown": ("markitdown",),
    "rapidocr_onnxruntime": ("rapidocr_onnxruntime", "rapidocr"),
    "onnxruntime": ("onnxruntime",),
    "win32com": ("pywin32", "pypiwin32"),
    "mss": ("mss",),
    "PIL": ("pillow",),
    "windows_use": ("windows_use", "windows-use"),
}


def now_stamp() -> str:
    return datetime.now().astimezone().strftime("%Y%m%d_%H%M%S")


def workspace_root() -> Path:
    return _WS_ROOT


def package_root() -> Path:
    return workspace_root().parent


def diagnostics_dir() -> Path:
    env = os.environ.get("LIG_DIAG_DIR", "").strip().strip('"')
    if env:
        return Path(env)
    return Path.home() / "OpenCodeLIG_USERDATA" / "diagnostics"


def safe_rel(path: Path) -> str:
    try:
        return str(path)
    except Exception:
        return repr(path)


def add(checks: list[Check], section: str, item: str, status: str,
        evidence: str, next_action: str = "") -> None:
    if status not in STATUSES:
        status = "WARN"
    checks.append(Check(section, item, status, evidence, next_action))


def has_module(name: str) -> bool:
    return importlib.util.find_spec(name) is not None


def run_cmd(args: list[str], timeout: int = 10, cwd: Path | None = None,
            env: dict[str, str] | None = None) -> dict[str, Any]:
    try:
        cp = subprocess.run(
            args,
            cwd=str(cwd or workspace_root()),
            env=env,
            capture_output=True,
            text=True,
            encoding="utf-8",
            errors="replace",
            timeout=timeout,
        )
        return {
            "ok": cp.returncode == 0,
            "returncode": cp.returncode,
            "stdout": (cp.stdout or "").strip()[-1500:],
            "stderr": (cp.stderr or "").strip()[-1500:],
        }
    except subprocess.TimeoutExpired as exc:
        return {
            "ok": False,
            "returncode": "TIMEOUT",
            "stdout": str(exc.stdout or "")[-800:],
            "stderr": str(exc.stderr or "")[-800:],
            "error": f"TIMEOUT after {timeout}s",
        }
    except Exception as exc:  # noqa: BLE001
        return {"ok": False, "error": f"{type(exc).__name__}: {exc}"[:800]}


def run_python_probe(code: str, timeout: int = 20, env: dict[str, str] | None = None) -> dict[str, Any]:
    merged = dict(os.environ)
    merged.update({"PYTHONUTF8": "1", "PYTHONIOENCODING": "utf-8"})
    if env:
        merged.update(env)
    return run_cmd([sys.executable, "-c", code], timeout=timeout, cwd=workspace_root(), env=merged)


def find_existing(paths: Iterable[Path]) -> list[str]:
    found = []
    for p in paths:
        try:
            if p and p.exists():
                found.append(str(p))
        except Exception:
            pass
    return found


def which_all(names: Iterable[str]) -> list[str]:
    found = []
    for name in names:
        hit = shutil.which(name)
        if hit and hit not in found:
            found.append(hit)
    return found


def common_program_paths() -> dict[str, list[Path]]:
    pf = Path(os.environ.get("ProgramFiles", r"C:\Program Files"))
    pfx86 = Path(os.environ.get("ProgramFiles(x86)", r"C:\Program Files (x86)"))
    local = Path(os.environ.get("LOCALAPPDATA", ""))
    return {
        "chrome": [
            pf / "Google" / "Chrome" / "Application" / "chrome.exe",
            pfx86 / "Google" / "Chrome" / "Application" / "chrome.exe",
        ],
        "obsidian": [
            local / "Programs" / "Obsidian" / "Obsidian.exe",
            pf / "Obsidian" / "Obsidian.exe",
        ],
        "matlab": [
            pf / "MATLAB" / "R2024a" / "bin" / "matlab.exe",
            pf / "MATLAB" / "R2023b" / "bin" / "matlab.exe",
        ],
        "autocad": [
            Path(r"C:\AutoCAD 2019\acad.exe"),
            Path(r"C:\AutoCAD 2019\accoreconsole.exe"),
            pf / "Autodesk" / "AutoCAD 2019" / "accoreconsole.exe",
            pf / "Autodesk" / "AutoCAD 2019" / "acad.exe",
            pf / "Autodesk" / "AutoCAD 2024" / "accoreconsole.exe",
            pf / "Autodesk" / "AutoCAD 2024" / "acad.exe",
        ],
        "fluent": [
            pf / "ANSYS Inc" / "v241" / "fluent" / "ntbin" / "win64" / "fluent.exe",
            pf / "ANSYS Inc" / "v242" / "fluent" / "ntbin" / "win64" / "fluent.exe",
        ],
        "solidworks": [
            pf / "SOLIDWORKS Corp" / "SOLIDWORKS" / "SLDWORKS.exe",
        ],
        "hwp": [
            pf / "Hnc" / "Office 2022" / "HOffice120" / "Bin" / "Hwp.exe",
            pfx86 / "Hnc" / "Office 2022" / "HOffice120" / "Bin" / "Hwp.exe",
        ],
    }


def registry_progid_exists(progid: str) -> bool:
    if platform.system().lower() != "windows":
        return False
    try:
        import winreg
        with winreg.OpenKey(winreg.HKEY_CLASSES_ROOT, progid):
            return True
    except Exception:
        return False


def sha256(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as fh:
        for chunk in iter(lambda: fh.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest()


def sanitize_evidence(text: Any, limit: int = 700) -> str:
    value = str(text).replace("|", "/").replace("\r", " ").replace("\n", " ")
    value = re.sub(r"Bearer\s+[A-Za-z0-9._\-+/=]+", "Bearer <hidden>", value)
    value = re.sub(r"(?i)(api[_-]?key\s*[=:]\s*)[^,\s;]+", r"\1<hidden>", value)
    return value[:limit]


def wheel_candidates(module: str) -> list[str]:
    wh = workspace_root() / "tools" / "wheelhouse"
    if not wh.exists():
        return []
    names = OPTIONAL_MODULE_WHEELS.get(module, (module,))
    out = []
    for wheel in wh.glob("*.whl"):
        lower = wheel.name.lower().replace("-", "_")
        if any(lower.startswith(name.lower().replace("-", "_")) for name in names):
            out.append(wheel.name)
    return sorted(out)


def check_core(checks: list[Check]) -> None:
    ws = workspace_root()
    root = package_root()
    py_launcher = which_all(["py.exe", "py"])
    powershell = run_cmd(["powershell", "-NoProfile", "-Command", "$PSVersionTable.PSVersion.ToString()"], timeout=5)
    add(checks, "필수 설치", "PC/OS", "PASS" if platform.system() == "Windows" else "WARN",
        f"user={os.environ.get('USERNAME','')}, computer={os.environ.get('COMPUTERNAME','')}, os={platform.platform()}, machine={platform.machine()}",
        "Windows PC에서 실행해야 실제 앱 자동화를 확인할 수 있습니다.")
    add(checks, "필수 설치", "PowerShell", "PASS" if powershell.get("ok") else "WARN",
        powershell.get("stdout") or powershell.get("stderr") or powershell.get("error", ""),
        "BAT 내부 진단 출력이 깨지면 PowerShell 실행 정책/인코딩을 확인하세요.")
    add(checks, "필수 설치", "Python 3.11", "PASS" if sys.version_info[:2] == (3, 11) else "FAIL",
        f"{sys.version.replace(chr(10), ' ')}; executable={sys.executable}; py_launcher={py_launcher}",
        "Python 3.11을 설치하고 py -3.11이 잡히는지 확인하세요.")
    add(checks, "필수 설치", "workspace", "PASS" if (ws / "agent_ops").is_dir() else "FAIL",
        safe_rel(ws), "설치 패키지를 다시 설치하세요.")
    add(checks, "필수 설치", "RUN_OPENCODE_LIG.bat", "PASS" if (ws / "RUN_OPENCODE_LIG.bat").exists() else "FAIL",
        safe_rel(ws / "RUN_OPENCODE_LIG.bat"), "workspace 복사 상태를 확인하세요.")

    payloads = [
        root / "payload" / "opencode.exe",
        Path.home() / "OpenCodeLIG" / "bin" / "opencode.exe",
    ]
    found = find_existing(payloads)
    add(checks, "필수 설치", "opencode.exe payload", "PASS" if found else "FAIL",
        "; ".join(found) if found else "not found",
        "패키지 payload/opencode.exe 또는 설치된 bin/opencode.exe가 필요합니다.")
    hashes = {}
    for p in payloads:
        if p.exists():
            try:
                hashes[str(p)] = sha256(p)
            except Exception as exc:  # noqa: BLE001
                hashes[str(p)] = f"hash_error:{exc!r}"
    if len(hashes) >= 2:
        values = list(hashes.values())
        add(checks, "필수 설치", "payload/installed hash match",
            "PASS" if len(set(values)) == 1 else "WARN",
            json.dumps(hashes, ensure_ascii=False),
            "패키지 payload와 설치된 bin이 다르면 설치 BAT를 다시 실행하세요.")
    for p in payloads:
        if p.exists():
            result = run_cmd([str(p), "--version"], timeout=10, cwd=ws)
            add(checks, "필수 설치", f"opencode 실행: {p.name}",
                "PASS" if result.get("ok") else "WARN",
                result.get("stdout") or result.get("stderr") or result.get("error", ""),
                "실행 실패 시 payload 손상 또는 보안 차단 여부를 확인하세요.")
            break

    sums = root / "SHA256SUMS.txt"
    payload = root / "payload" / "opencode.exe"
    if sums.exists() and payload.exists():
        text = sums.read_text(encoding="utf-8", errors="replace").lower()
        actual = sha256(payload)
        add(checks, "필수 설치", "payload SHA256", "PASS" if actual in text else "FAIL",
            actual, "SHA256SUMS.txt와 payload가 같은 패키지에서 온 것인지 확인하세요.")

    path_opencode = shutil.which("opencode")
    add(checks, "필수 설치", "opencode PATH", "PASS" if path_opencode else "WARN",
        path_opencode or "PATH에는 없음",
        "일반 셸에서 opencode를 직접 칠 필요가 있으면 %USERPROFILE%\\OpenCodeLIG\\bin을 PATH에 추가하세요.")


def check_opencode_patch(checks: list[Check]) -> None:
    exe = package_root() / "payload" / "opencode.exe"
    if not exe.exists():
        exe = Path.home() / "OpenCodeLIG" / "bin" / "opencode.exe"
    if not exe.exists():
        add(checks, "OpenCode 패치", "permission badge payload", "FAIL", "opencode.exe not found",
            "payload 또는 설치 bin을 확인하세요.")
        return
    blob = exe.read_bytes()
    tokens = {token: (token.encode("utf-8") in blob) for token in ["ASK", "AUTO", "FULL", "Permission approval", "/permission"]}
    add(checks, "OpenCode 패치", "ASK/AUTO/FULL 표시 문자열", "PASS" if all(tokens.values()) else "WARN",
        json.dumps(tokens, ensure_ascii=False),
        "ASK/AUTO/FULL 모드 표시 패치가 payload에 포함됐는지 확인하세요.")
    for args in (["--help"], ["run", "--help"], ["models", "--help"]):
        res = run_cmd([str(exe), *args], timeout=12, cwd=workspace_root())
        add(checks, "OpenCode 패치", "opencode " + " ".join(args),
            "PASS" if res.get("ok") else "WARN",
            res.get("stdout") or res.get("stderr") or res.get("error", ""),
            "기본 명령 help가 실패하면 빌드 산출물/런타임 DLL 차단 여부를 확인하세요.")


def check_config_and_loop(checks: list[Check]) -> None:
    ws = workspace_root()
    try:
        from agent_ops.lig_providers import SECRET_ENV_PATH, build_providers, load_lig_env, validate_config
        env = load_lig_env()
        cfg = validate_config()
        secret_meta = {
            "secret_file": str(SECRET_ENV_PATH),
            "secret_exists": SECRET_ENV_PATH.exists(),
            "key_present": bool(env.get("LIG_API_KEY")),
            "key_length": len(env.get("LIG_API_KEY", "")) if env.get("LIG_API_KEY") else 0,
            "url_present": bool(env.get("LIG_GATEWAY_BASE_URL") or env.get("LIG_LOCAL_BASE_URL")),
            "profile": cfg.get("profile"),
        }
        add(checks, "모델/설정", "게이트웨이 secret 형식", "PASS" if cfg.get("ready") else "FAIL",
            json.dumps(secret_meta, ensure_ascii=False),
            "USERDATA\\secrets\\lig-api.env를 확인하세요. URL/API 키 값은 보고서에 숨깁니다.")
        providers = build_providers()
        ready_routes = [name for name, p in providers.items() if p.get("base_url") and p.get("model")]
        add(checks, "모델/설정", "provider routes", "PASS" if ready_routes else "FAIL",
            ", ".join(ready_routes), "config/lig-api.env.example 및 USERDATA secret을 확인하세요.")
        route_results = []
        ok_count = 0
        for name, p in providers.items():
            base = str(p.get("base_url", "")).rstrip("/")
            if not base:
                route_results.append({"route": name, "status": "missing_base"})
                continue
            try:
                req = urllib.request.Request(
                    base + "/models",
                    headers={"Authorization": "Bearer " + str(env.get("LIG_API_KEY", ""))},
                )
                with urllib.request.urlopen(req, timeout=5) as resp:
                    ok = 200 <= int(resp.status) < 500
                    ok_count += 1 if ok else 0
                    route_results.append({"route": name, "http": int(resp.status), "model": p.get("model")})
            except Exception as exc:  # noqa: BLE001
                route_results.append({"route": name, "error": type(exc).__name__, "model": p.get("model")})
        add(checks, "모델/설정", "게이트웨이 /models live probe",
            "PASS" if ok_count else "PENDING",
            json.dumps(route_results, ensure_ascii=False),
            "사내망에서 실행했는데 reachable 0이면 게이트웨이/프록시/키를 확인하세요.")
    except Exception as exc:  # noqa: BLE001
        add(checks, "모델/설정", "게이트웨이 설정", "FAIL", repr(exc)[:500],
            "agent_ops.lig_providers import 오류를 확인하세요.")

    opencode_json = ws / "opencode.json"
    try:
        data = json.loads(opencode_json.read_text(encoding="utf-8-sig"))
        model = str(data.get("model", ""))
        providers_cfg = sorted((data.get("provider") or {}).keys())
        add(checks, "모델/설정", "opencode.json", "PASS" if "/" in model else "FAIL",
            f"default_model={model}, providers={providers_cfg}", "opencode.json JSON 및 provider/model 형식을 확인하세요.")
    except Exception as exc:  # noqa: BLE001
        add(checks, "모델/설정", "opencode.json", "FAIL", repr(exc)[:500], "opencode.json을 복구하세요.")

    try:
        from agent_ops.intelligence_map import all_items, coverage_summary
        items = all_items()
        pending = [asdict(i) for i in items if i.status == "pending"]
        auto = [asdict(i) for i in items if i.status == "auto"]
        add(checks, "자동지능 루프", "intelligence map", "PASS",
            f"items={len(items)}, auto={len(auto)}, pending={len(pending)}, summary={json.dumps(coverage_summary(), ensure_ascii=False, sort_keys=True)}")
    except Exception as exc:  # noqa: BLE001
        add(checks, "자동지능 루프", "intelligence map", "FAIL", repr(exc)[:500],
            "agent_ops.intelligence_map 오류를 확인하세요.")

    tasks = [
        ("browser", "크롬으로 회사 포털 확인해줘"),
        ("memory", "이 내용 기억해: 사내 점검은 한번에 끝내야 한다"),
        ("wiki", "위키 정리해줘"),
        ("document", "회의록을 워드 보고서로 정리해줘"),
        ("engineering", "알루미늄 피로 설계 기준을 교수처럼 설명해줘"),
        ("danger", "C 드라이브의 모든 파일을 삭제해줘"),
    ]
    with tempfile.TemporaryDirectory(prefix="agentops_pending_") as td:
        env = dict(os.environ)
        env.update({
            "PYTHONUTF8": "1",
            "PYTHONIOENCODING": "utf-8",
            "AGENTOPS_ROOT": str(Path(td) / "workspace"),
            "AGENTOPS_MEMORY_DIR": str(Path(td) / "memory"),
            "LIG_DIAG_DIR": str(Path(td) / "diag"),
        })
        matrix = []
        ok_rows = 0
        for label, task in tasks:
            task_file = Path(td) / f"task_{label}.txt"
            task_file.write_text(task, encoding="utf-8")
            cmd = [sys.executable, str(ws / "agent_ops" / "agentops.py"), "auto", "--task-file", str(task_file), "--dry-run"]
            cp = subprocess.run(cmd, cwd=str(ws), env=env, capture_output=True,
                                text=True, encoding="utf-8", errors="replace", timeout=30)
            trace = Path(td) / "diag" / "auto-route-last.json"
            trace_data = json.loads(trace.read_text(encoding="utf-8")) if trace.exists() else {}
            has_fields = all(k in trace_data for k in ("capability_ids", "policy", "route_hints", "evaluation", "memory_hooks"))
            ok_rows += 1 if cp.returncode == 0 and has_fields else 0
            matrix.append({
                "case": label,
                "returncode": cp.returncode,
                "command": trace_data.get("command"),
                "path": trace_data.get("selected_path"),
                "policy": (trace_data.get("policy") or {}).get("mode"),
                "caps": trace_data.get("capability_ids", [])[:3],
                "blocked": (trace_data.get("policy") or {}).get("blocked"),
            })
        matrix_summary = "; ".join(
            f"{row['case']}=>rc{row['returncode']}/{row['command']}/{row['path']}/policy={row['policy']}/blocked={row['blocked']}/caps={','.join(row['caps'])}"
            for row in matrix
        )
        add(checks, "자동지능 루프", "/auto dry-run 시나리오 매트릭스",
            "PASS" if ok_rows == len(tasks) else "FAIL",
            matrix_summary,
            "각 케이스에 capability/policy/route/evaluation/memory hook이 기록되어야 합니다.")


def check_optional_modules(checks: list[Check]) -> None:
    modules = {
        "openpyxl": "xlsx 생성/읽기",
        "docx": "Word docx 생성",
        "pptx": "PowerPoint pptx 생성",
        "markitdown": "PDF/Office 문서 읽기",
        "rapidocr_onnxruntime": "OCR",
        "onnxruntime": "OCR backend",
        "win32com": "Office/HWP/Outlook/SolidWorks COM",
        "mss": "화면 캡처",
        "PIL": "이미지 처리",
        "windows_use": "Desktop UI 자동화",
    }
    module_state = {}
    for mod, purpose in modules.items():
        installed = has_module(mod)
        candidates = wheel_candidates(mod)
        module_state[mod] = {"installed": installed, "wheel_candidates": candidates}
        status = "PASS" if installed else ("PENDING" if candidates else "WARN")
        next_action = "이미 import 가능합니다." if installed else (
            "wheelhouse 후보가 있으므로 install-tools.bat로 오프라인 설치 가능성을 확인하세요."
            if candidates else "필요 기능이면 해당 wheel을 tools\\wheelhouse에 반입하세요."
        )
        add(checks, "오프라인 의존성", mod, status,
            f"{purpose}; wheels={candidates}", next_action)
    wh = workspace_root() / "tools" / "wheelhouse"
    wheels = list(wh.glob("*.whl")) if wh.exists() else []
    add(checks, "오프라인 의존성", "wheelhouse", "PASS" if wheels else "WARN",
        f"{len(wheels)} wheel(s) at {wh}", "오프라인 추가 설치가 필요하면 wheelhouse를 채우세요.")

    smoke_code = r"""
import json, tempfile
from pathlib import Path
out = {}
try:
    import openpyxl
    p = Path(tempfile.mkdtemp()) / "probe.xlsx"
    wb = openpyxl.Workbook(); wb.active["A1"] = 42; wb.save(p)
    out["openpyxl"] = {"ok": p.exists(), "size": p.stat().st_size}
except Exception as exc:
    out["openpyxl"] = {"ok": False, "error": type(exc).__name__}
try:
    import docx
    p = Path(tempfile.mkdtemp()) / "probe.docx"
    d = docx.Document(); d.add_paragraph("probe"); d.save(p)
    out["docx"] = {"ok": p.exists(), "size": p.stat().st_size}
except Exception as exc:
    out["docx"] = {"ok": False, "error": type(exc).__name__}
try:
    import pptx
    p = Path(tempfile.mkdtemp()) / "probe.pptx"
    prs = pptx.Presentation(); prs.slides.add_slide(prs.slide_layouts[0]); prs.save(p)
    out["pptx"] = {"ok": p.exists(), "size": p.stat().st_size}
except Exception as exc:
    out["pptx"] = {"ok": False, "error": type(exc).__name__}
print(json.dumps(out, ensure_ascii=False))
"""
    smoke = run_python_probe(smoke_code, timeout=25)
    add(checks, "오프라인 의존성", "문서 생성 라이브러리 smoke",
        "PASS" if smoke.get("ok") else "PENDING",
        smoke.get("stdout") or smoke.get("stderr") or smoke.get("error", ""),
        "docx/xlsx/pptx 생성 실패 항목은 해당 wheel 설치 상태를 확인하세요.")


def com_activation_probe(progid: str) -> dict[str, Any]:
    code = f"""
import json
out={{"progid": {progid!r}}}
try:
    import win32com.client
    app = win32com.client.Dispatch({progid!r})
    out["ok"] = True
    out["type"] = str(type(app))
    try:
        app.Visible = False
    except Exception:
        pass
    try:
        app.Quit()
    except Exception:
        pass
except Exception as exc:
    out["ok"] = False
    out["error"] = type(exc).__name__
    out["detail"] = str(exc)[:200]
print(json.dumps(out, ensure_ascii=False))
"""
    return run_python_probe(code, timeout=20)


def check_apps_and_pending(checks: list[Check]) -> None:
    paths = common_program_paths()
    env_paths = {
        "matlab": os.environ.get("MATLAB_EXE", ""),
        "autocad": os.environ.get("ACCORECONSOLE_EXE", "") or os.environ.get("ACAD_EXE", "") or os.environ.get("AUTOCAD_EXE", ""),
        "fluent": os.environ.get("FLUENT_EXE", ""),
    }
    env_names = {
        "matlab": "MATLAB_EXE",
        "autocad": "ACCORECONSOLE_EXE/ACAD_EXE",
        "fluent": "FLUENT_EXE",
    }
    path_names = {
        "matlab": ["matlab.exe", "matlab"],
        "autocad": ["accoreconsole.exe", "accoreconsole", "acad.exe", "acad"],
        "fluent": ["fluent.exe", "fluent"],
    }
    executable_hits: dict[str, list[str]] = {}
    for key, env_path in env_paths.items():
        candidates = ([Path(env_path)] if env_path else []) + paths.get(key, [])
        found = find_existing(candidates) + which_all(path_names[key])
        executable_hits[key] = found
        add(checks, "앱/도구 pending", f"{key} executable", "PASS" if found else "PENDING",
            "; ".join(found) if found else "not found",
            f"{env_names[key]} 환경변수, PATH, 또는 표준 설치 경로를 확인하세요.")

    cli_probes = {
        "autocad": (["/?"], 20),
        "fluent": (["-help"], 20),
    }
    for key, (extra, timeout) in cli_probes.items():
        if not executable_hits.get(key):
            continue
        if key == "autocad" and Path(executable_hits[key][0]).name.lower() == "acad.exe":
            add(checks, "앱/도구 pending", f"{key} cli probe", "SKIP",
                f"{executable_hits[key][0]} is GUI AutoCAD; skipped /? probe to avoid launching UI",
                "실제 .scr 실행은 사용자 작업 요청 때 사본 DWG + /p LIGNEX1 /product ACADM /b 기준으로 수행합니다.")
            continue
        res = run_cmd([executable_hits[key][0], *extra], timeout=timeout)
        add(checks, "앱/도구 pending", f"{key} cli probe",
            "PASS" if res.get("ok") else "WARN",
            res.get("stdout") or res.get("stderr") or res.get("error", ""),
            "help/version probe만 실행했습니다. 실제 파일 실행은 사용자 작업 요청 때 사본/임시파일 기준으로 수행합니다.")

    chrome_found = find_existing(paths["chrome"]) + which_all(["chrome.exe", "chrome"])
    add(checks, "앱/도구 pending", "Chrome", "PASS" if chrome_found else "PENDING",
        "; ".join(chrome_found) if chrome_found else "not found",
        "Chrome 설치 및 launch\\chrome-debug.bat 실행 여부를 확인하세요.")
    try:
        with urllib.request.urlopen("http://127.0.0.1:9222/json/version", timeout=2) as resp:
            version = json.loads(resp.read().decode("utf-8", errors="replace"))
        with urllib.request.urlopen("http://127.0.0.1:9222/json/list", timeout=2) as resp:
            tabs = json.loads(resp.read().decode("utf-8", errors="replace"))
        cdp_evidence = {"browser": version.get("Browser"), "webSocketDebuggerUrl": bool(version.get("webSocketDebuggerUrl")), "tabs": len(tabs)}
        cdp_ok = True
    except Exception as exc:  # noqa: BLE001
        cdp_evidence = {"reachable": False, "error": type(exc).__name__}
        cdp_ok = False
    add(checks, "앱/도구 pending", "Chrome CDP 9222", "PASS" if cdp_ok else "PENDING",
        json.dumps(cdp_evidence, ensure_ascii=False),
        "브라우저 자동화 검증 전 launch\\chrome-debug.bat로 디버그 크롬을 여세요.")

    progids = {
        "Excel COM": "Excel.Application",
        "Word COM": "Word.Application",
        "PowerPoint COM": "PowerPoint.Application",
        "Outlook COM": "Outlook.Application",
        "HWP COM": "HWPFrame.HwpObject",
        "SolidWorks COM": "SldWorks.Application",
    }
    for label, progid in progids.items():
        exists = registry_progid_exists(progid)
        add(checks, "앱/도구 pending", f"{label} registry", "PASS" if exists else "PENDING",
            f"ProgID={progid}", f"{label} 등록/설치 상태를 확인하세요.")
        if exists and has_module("win32com"):
            res = com_activation_probe(progid)
            add(checks, "앱/도구 pending", f"{label} activation",
                "PASS" if res.get("ok") and '"ok": true' in (res.get("stdout") or "").lower() else "WARN",
                res.get("stdout") or res.get("stderr") or res.get("error", ""),
                "COM 객체 생성 단계입니다. 실제 문서/모델 파일 조작은 사본 정책으로 별도 수행됩니다.")

    try:
        from agent_ops.adapters import ADAPTERS
        for adapter_id, spec in ADAPTERS.items():
            pending = str(spec.get("pending", "") or "")
            available = bool(spec.get("available"))
            validated = str(spec.get("validated", "") or spec.get("home_smoke", "") or "")
            if pending:
                status = "WARN" if available else "PENDING"
                evidence = f"available={available}, validated={'yes' if validated else 'no'}, pending={pending}"
                add(checks, "앱/도구 pending", f"adapter:{adapter_id}", status, evidence,
                    "위 항목의 실행파일/COM/의존성 결과와 함께 판단하세요.")
    except Exception as exc:  # noqa: BLE001
        add(checks, "앱/도구 pending", "adapter inventory", "FAIL", repr(exc)[:500],
            "agent_ops.adapters import 오류를 확인하세요.")


def check_screen_ocr(checks: list[Check]) -> None:
    if has_module("mss"):
        code = r"""
import json
try:
    import mss
    with mss.mss() as sct:
        mon = sct.monitors[0]
        img = sct.grab({"left": mon["left"], "top": mon["top"], "width": min(80, mon["width"]), "height": min(80, mon["height"])})
    print(json.dumps({"ok": True, "size": [img.width, img.height], "monitors": len(sct.monitors) if hasattr(sct, "monitors") else None}, ensure_ascii=False))
except Exception as exc:
    print(json.dumps({"ok": False, "error": type(exc).__name__, "detail": str(exc)[:200]}, ensure_ascii=False))
    raise SystemExit(1)
"""
        res = run_python_probe(code, timeout=15)
        add(checks, "화면/OCR", "mss screenshot smoke", "PASS" if res.get("ok") else "WARN",
            res.get("stdout") or res.get("stderr") or res.get("error", ""),
            "화면 캡처 권한 또는 보안 프로그램 차단 여부를 확인하세요.")
    else:
        add(checks, "화면/OCR", "mss screenshot smoke", "PENDING", "mss module missing",
            "화면 인식 자동화가 필요하면 mss/Pillow wheel을 반입하세요.")

    if has_module("rapidocr_onnxruntime"):
        code = r"""
import json
try:
    from rapidocr_onnxruntime import RapidOCR
    engine = RapidOCR()
    print(json.dumps({"ok": True, "engine": str(type(engine))}, ensure_ascii=False))
except Exception as exc:
    print(json.dumps({"ok": False, "error": type(exc).__name__, "detail": str(exc)[:200]}, ensure_ascii=False))
    raise SystemExit(1)
"""
        res = run_python_probe(code, timeout=30)
        add(checks, "화면/OCR", "RapidOCR instantiate", "PASS" if res.get("ok") else "WARN",
            res.get("stdout") or res.get("stderr") or res.get("error", ""),
            "OCR 모델 파일/onnxruntime wheel 호환성을 확인하세요.")
    else:
        add(checks, "화면/OCR", "RapidOCR instantiate", "PENDING", "rapidocr_onnxruntime missing",
            "OCR이 필요하면 RapidOCR/onnxruntime wheel 및 모델을 반입하세요.")


def check_obsidian_wiki(checks: list[Check]) -> None:
    paths = common_program_paths()
    obsidian_found = find_existing(paths["obsidian"]) + which_all(["obsidian.exe", "obsidian"])
    wiki_dir = Path.home() / "OpenCodeLIG_USERDATA" / "memory" / "wiki"
    add(checks, "Obsidian/위키", "Obsidian executable", "PASS" if obsidian_found else "PENDING",
        "; ".join(obsidian_found) if obsidian_found else "not found",
        "Obsidian을 설치하고 USERDATA\\memory\\wiki를 vault로 열 수 있게 준비하세요.")
    add(checks, "Obsidian/위키", "실 USERDATA wiki vault", "PASS" if wiki_dir.exists() else "WARN",
        f"{wiki_dir}; obsidian_config={(wiki_dir / '.obsidian').exists()}",
        "첫 기억/위키 실행 후 자동 생성됩니다. 직접 쓰는 노트는 wiki\\manual에 둡니다.")

    with tempfile.TemporaryDirectory(prefix="agentops_wiki_probe_") as td:
        env = dict(os.environ)
        env.update({
            "PYTHONUTF8": "1",
            "PYTHONIOENCODING": "utf-8",
            "AGENTOPS_ROOT": str(Path(td) / "workspace"),
            "AGENTOPS_MEMORY_DIR": str(Path(td) / "memory"),
            "LIG_DIAG_DIR": str(Path(td) / "diag"),
        })
        cmd1 = [sys.executable, str(workspace_root() / "agent_ops" / "agentops.py"), "remember", "--title", "pending-check", "위키 점검용 격리 메모리"]
        r1 = run_cmd(cmd1, timeout=20, cwd=workspace_root(), env=env)
        cmd2 = [sys.executable, str(workspace_root() / "agent_ops" / "agentops.py"), "wiki"]
        r2 = run_cmd(cmd2, timeout=30, cwd=workspace_root(), env=env)
        memory = Path(td) / "memory"
        generated = sorted(str(p.relative_to(memory)) for p in memory.rglob("*.md")) if memory.exists() else []
        add(checks, "Obsidian/위키", "격리 remember→wiki smoke",
            "PASS" if r1.get("ok") and r2.get("ok") and generated else "FAIL",
            json.dumps({"remember_rc": r1.get("returncode"), "wiki_rc": r2.get("returncode"), "generated": generated[:8]}, ensure_ascii=False),
            "격리 smoke 실패 시 memory_manager/wiki_manager 경로를 확인하세요. 실제 USERDATA는 건드리지 않았습니다.")

    autosave = workspace_root() / ".opencode" / "plugins" / "session-autosave.ts"
    autosave_text = autosave.read_text(encoding="utf-8", errors="replace") if autosave.exists() else ""
    autosave_ok = (
        autosave.exists()
        and "wiki\", \"sessions\"" in autosave_text
        and "appendFileSync(sessionFile()" in autosave_text
        and "log-activity" in autosave_text
        and "(?i:" not in autosave_text
    )
    autosave_sessions = "wiki\", \"sessions\"" in autosave_text
    autosave_promote = "log-activity" in autosave_text
    add(checks, "Obsidian/위키", "세션 자동저장 플러그인", "PASS" if autosave_ok else "WARN",
        f"{autosave}; exists={autosave.exists()}; sessions={autosave_sessions}; memory_promote={autosave_promote}",
        "대화 중 창을 닫아도 Obsidian wiki\\sessions 노트에 남고, 일부는 저우선순위 기억으로 승격되어야 합니다.")


def check_command_surfaces(checks: list[Check]) -> None:
    """Smoke-test user-visible command surfaces in an isolated workspace.

    This complements /auto routing: if a user asks naturally, /auto chooses the
    route; this probe confirms the target command behind that route can at least
    start, parse arguments, and write its expected local artifact without touching
    real USERDATA.
    """
    ws = workspace_root()
    commands_dir = ws / ".opencode" / "commands"
    agents_dir = ws / ".opencode" / "agents"
    plugins_dir = ws / ".opencode" / "plugins"
    command_files = sorted(p.name for p in commands_dir.glob("*.md")) if commands_dir.exists() else []
    add(checks, "명령/산출물 smoke", ".opencode command docs",
        "PASS" if len(command_files) >= 30 else "WARN",
        f"count={len(command_files)}, sample={command_files[:8]}",
        "OpenCode slash command 문서가 누락되면 TUI에서 명령 안내가 약해집니다.")
    agent_files = sorted(p.name for p in agents_dir.glob("*.md")) if agents_dir.exists() else []
    plugin_files = sorted(p.name for p in plugins_dir.glob("*.ts")) if plugins_dir.exists() else []
    add(checks, "명령/산출물 smoke", ".opencode agents/plugins",
        "PASS" if agent_files and plugin_files else "FAIL",
        f"agents={len(agent_files)}, plugins={plugin_files}",
        "에이전트/플러그인 파일이 없으면 기억 주입, command guard, 상태 연동이 빠집니다.")

    with tempfile.TemporaryDirectory(prefix="agentops_cmd_probe_") as td:
        td_path = Path(td)
        out_dir = td_path / "out"
        out_dir.mkdir(parents=True, exist_ok=True)
        csv = td_path / "probe.csv"
        csv.write_text("name,value\nalpha,10\nbeta,20\n", encoding="utf-8")
        doc_spec = td_path / "doc.json"
        doc_spec.write_text(json.dumps({
            "title": "pending-check-doc",
            "sections": [{"heading": "요약", "paragraphs": ["점검 문서"], "bullets": ["A", "B"]}],
        }, ensure_ascii=False), encoding="utf-8")
        ppt_spec = td_path / "ppt.json"
        ppt_spec.write_text(json.dumps({
            "title": "pending-check-ppt",
            "slides": [{"title": "핵심", "points": ["점검", "완료"]}],
        }, ensure_ascii=False), encoding="utf-8")
        task_file = td_path / "work-task.txt"
        task_file.write_text("간단한 점검 산출물 만들어줘", encoding="utf-8")
        isolated_root = td_path / "workspace"
        isolated_root.mkdir(parents=True, exist_ok=True)
        safe_probe_dir = ws / "agent_ops" / "results" / "pending_check"
        safe_probe_dir.mkdir(parents=True, exist_ok=True)
        safe_content = safe_probe_dir / "safe_write_content.txt"
        safe_content.write_text("safe-write probe", encoding="utf-8")

        env = dict(os.environ)
        env.update({
            "PYTHONUTF8": "1",
            "PYTHONIOENCODING": "utf-8",
            "AGENTOPS_ROOT": str(isolated_root),
            "AGENTOPS_MEMORY_DIR": str(td_path / "memory"),
            "LIG_DIAG_DIR": str(td_path / "diag"),
            "AGENTOPS_OUTPUT_DIR": str(out_dir),
        })
        agentops = str(ws / "agent_ops" / "agentops.py")
        probes: list[tuple[str, list[str], set[int], str]] = [
            ("doctor", [agentops, "doctor"], {0}, "상태 진단"),
            ("deps", [agentops, "deps"], {0}, "선택 의존성 표시"),
            ("status", [agentops, "status"], {0}, "상태 JSON"),
            ("report", [agentops, "report"], {0}, "운영 보고서"),
            ("memorycheck", [agentops, "memorycheck"], {0}, "기억 상태"),
            ("recall-pinned", [agentops, "recall", "--pinned"], {0}, "세션 시작 기억 주입"),
            ("routine-list", [agentops, "routine", "list"], {0}, "루틴 목록"),
            ("schedule-list", [agentops, "schedule", "list", "--when", "all"], {0}, "일정 목록"),
            ("schedule-today", [agentops, "schedule", "today"], {0}, "오늘 일정"),
            ("timeline", [agentops, "timeline", "--gap", "1"], {0}, "활동 타임라인"),
            ("briefing", [agentops, "briefing"], {0}, "브리핑"),
            ("weekly", [agentops, "weekly"], {0}, "주간 보고"),
            ("book", [agentops, "book"], {0}, "지식책"),
            ("watch", [agentops, "watch", "--max-age", "999999"], {0, 3, 4}, "무한대기 감시 판정"),
            ("checkpoint", [agentops, "checkpoint", "--note", "pending-check"], {0}, "체크포인트"),
            ("continue-once", [agentops, "continue-once"], {0}, "이어가기 상태"),
            ("enqueue", [agentops, "enqueue", "pending-check task", "--kind", "doctor"], {0}, "큐 등록"),
            ("safety-check", [agentops, "safety-check", "Remove-Item -Recurse C:\\"], {0}, "위험 명령 분류"),
            ("report-html", [agentops, "report-html", "--input", str(csv), "--title", "probe"], {0}, "HTML 리포트"),
            ("report-xlsx", [agentops, "report-xlsx", "--input", str(csv), "--out", str(out_dir / "probe.xlsx")], {0}, "XLSX 리포트"),
            ("office-docx", [agentops, "office-doc", "--kind", "docx", "--spec", str(doc_spec), "--out", str(out_dir / "probe.docx")], {0}, "DOCX 생성"),
            ("office-pptx", [agentops, "office-doc", "--kind", "pptx", "--spec", str(ppt_spec), "--out", str(out_dir / "probe.pptx")], {0}, "PPTX 생성"),
            ("doc-template", [agentops, "doc-template", "회의록", "--out", str(out_dir), "--title", "probe"], {0}, "정형문서 생성"),
            ("safe-write", [agentops, "safe-write",
                            "agent_ops/results/pending_check/safe_write_target.txt",
                            "agent_ops/results/pending_check/safe_write_content.txt"], {0}, "안전 파일 생성"),
            ("work-mock", [agentops, "work", "--task-file", str(task_file), "--mode", "mock"], {0}, "작업 산출물 mock"),
        ]
        rows = []
        ok_count = 0
        for label, cmd, accepted, purpose in probes:
            res = run_cmd([sys.executable, *cmd], timeout=35, cwd=ws, env=env)
            rc = res.get("returncode")
            ok = rc in accepted
            ok_count += 1 if ok else 0
            rows.append(f"{label}=rc{rc}:{'ok' if ok else 'bad'}")
            status = "PASS" if ok else "WARN"
            add(checks, "명령/산출물 smoke", f"command:{label}", status,
                f"{purpose}; rc={rc}; out={(res.get('stdout') or res.get('stderr') or res.get('error', ''))[:180]}",
                "실패 시 해당 명령의 인자 파싱/의존성/격리 USERDATA 쓰기 경로를 확인하세요.")
        generated = sorted(str(p.relative_to(out_dir)) for p in out_dir.rglob("*") if p.is_file())
        expected_any = {"probe.xlsx", "probe.docx", "probe.pptx"}
        generated_names = {Path(p).name for p in generated}
        add(checks, "명령/산출물 smoke", "isolated artifact outputs",
            "PASS" if expected_any.issubset(generated_names) else "WARN",
            f"commands_ok={ok_count}/{len(probes)}, generated={generated[:12]}, summary={'; '.join(rows)}",
            "문서/표/슬라이드 산출물이 없으면 Python 문서 생성 의존성 또는 output dir 정책을 확인하세요.")


def check_docs_package(checks: list[Check]) -> None:
    ws = workspace_root()
    required_docs = [
        ws / "docs" / "사용법" / "GUIDE.md",
        ws / "docs" / "사용법" / "RUNBOOK.md",
        ws / "docs" / "운영" / "INTELLIGENCE_COVERAGE_REPORT.md",
    ]
    for p in required_docs:
        text = p.read_text(encoding="utf-8", errors="replace") if p.exists() else ""
        add(checks, "문서/패키지", p.name, "PASS" if p.exists() and len(text) > 500 else "FAIL",
            f"{safe_rel(p)}; chars={len(text)}", "사용자용 패키지에 충분한 문서가 포함되어야 합니다.")
    tests_dir = ws / "tests"
    add(checks, "문서/패키지", "사용자 패키지 테스트폴더", "WARN" if tests_dir.exists() else "PASS",
        f"tests_dir_exists={tests_dir.exists()}",
        "배포 패키지에서는 tests 폴더가 제외되는 것이 정상입니다.")
    for p in [package_root() / "점검용_전체확인.bat", ws / "점검용_전체확인.bat"]:
        if p.exists():
            raw = p.read_bytes()
            crlf_ok = b"\n" not in raw.replace(b"\r\n", b"")
            chcp_ok = b"chcp 65001" in raw
            ok = chcp_ok and crlf_ok
            add(checks, "문서/패키지", f"점검 BAT: {p.name}", "PASS" if ok else "FAIL",
                f"{p}; crlf={crlf_ok}; chcp={chcp_ok}",
                ".bat는 CRLF + chcp 65001이어야 합니다.")
        else:
            add(checks, "문서/패키지", f"점검 BAT: {p.name}", "WARN", str(p), "패키지 루트와 설치 workspace에 점검 BAT가 있어야 합니다.")
    guard = ws / ".opencode" / "plugins" / "command-guard.ts"
    guard_text = guard.read_text(encoding="utf-8", errors="replace") if guard.exists() else ""
    deny_hits = [term for term in ["reset --hard", "Remove-Item", "format", "rmdir"] if term.lower() in guard_text.lower()]
    add(checks, "문서/패키지", "command guard plugin", "PASS" if guard.exists() and deny_hits else "FAIL",
        f"{guard}; deny_hits={deny_hits}",
        "ASK/AUTO/FULL 어느 모드에서도 위험 명령 가드는 유지되어야 합니다.")


def build_report(checks: list[Check], report_id: str) -> tuple[dict[str, Any], str]:
    counts = {s: sum(1 for c in checks if c.status == s) for s in STATUSES}
    blocking = [c for c in checks if c.status == "FAIL"]
    pending = [c for c in checks if c.status in {"PENDING", "WARN"}]
    payload = {
        "report_id": report_id,
        "timestamp": datetime.now().astimezone().isoformat(timespec="seconds"),
        "workspace": str(workspace_root()),
        "counts": counts,
        "blocking_failures": [asdict(c) for c in blocking],
        "pending_or_warning": [asdict(c) for c in pending],
        "checks": [asdict(c) for c in checks],
    }
    verdict = "핵심 실패 없음" if not blocking else "핵심 실패 있음"
    lines = [
        "# OpenCodeLIG 사내 PC 미결항목 통합 점검 보고서",
        "",
        "## 보내줄 파일",
        "",
        f"- 이 파일 하나만 보내면 됩니다: `{report_id}.md` 또는 같은 폴더의 `pending-check-last.md`",
        "- API 키와 게이트웨이 URL은 의도적으로 숨겼습니다.",
        "",
        "## 요약",
        "",
        f"- report_id: `{report_id}`",
        f"- timestamp: `{payload['timestamp']}`",
        f"- workspace: `{payload['workspace']}`",
        f"- verdict: **{verdict}**",
        f"- counts: PASS {counts['PASS']} / WARN {counts['WARN']} / PENDING {counts['PENDING']} / FAIL {counts['FAIL']} / SKIP {counts['SKIP']}",
        "",
        "## 판정 기준",
        "",
        "- FAIL: 필수 설치/핵심 자동 루프가 깨진 항목입니다. 먼저 해결해야 합니다.",
        "- PENDING: 사내 앱, COM, 브라우저, OCR, Fluent 등 현장 검증이 필요한 항목입니다.",
        "- WARN: 동작은 가능하지만 직접 실행 경로, 문서, 선택 기능 확인이 남은 항목입니다.",
        "- PASS: 현재 PC에서 확인된 항목입니다.",
        "",
    ]
    if blocking:
        lines += ["## 먼저 봐야 할 FAIL", "", "| section | item | evidence | next action |", "|---|---|---|---|"]
        for c in blocking:
            lines.append(f"| {sanitize_evidence(c.section)} | {sanitize_evidence(c.item)} | {sanitize_evidence(c.evidence)} | {sanitize_evidence(c.next_action, 240)} |")
        lines.append("")
    if pending:
        lines += ["## 미결/WARN 빠른 목록", "", "| status | section | item | next action |", "|---|---|---|---|"]
        for c in pending:
            lines.append(f"| {c.status} | {sanitize_evidence(c.section)} | {sanitize_evidence(c.item)} | {sanitize_evidence(c.next_action, 240)} |")
        lines.append("")
    for section in sorted({c.section for c in checks}):
        lines += [f"## {section}", "", "| status | item | evidence | next action |", "|---|---|---|---|"]
        for c in [x for x in checks if x.section == section]:
            lines.append(f"| {c.status} | {sanitize_evidence(c.item)} | {sanitize_evidence(c.evidence)} | {sanitize_evidence(c.next_action, 240)} |")
        lines.append("")
    return payload, "\n".join(lines)


def main(argv: list[str] | None = None) -> int:
    try:
        sys.stdout.reconfigure(encoding="utf-8", errors="replace")
        sys.stderr.reconfigure(encoding="utf-8", errors="replace")
    except Exception:
        pass
    parser = argparse.ArgumentParser()
    parser.add_argument("--out-dir", default="")
    args = parser.parse_args(argv)
    out_dir = Path(args.out_dir) if args.out_dir else diagnostics_dir() / "pending_checks"
    out_dir.mkdir(parents=True, exist_ok=True)
    report_id = "pending_check_" + now_stamp()

    checks: list[Check] = []
    check_core(checks)
    check_opencode_patch(checks)
    check_config_and_loop(checks)
    check_optional_modules(checks)
    check_apps_and_pending(checks)
    check_screen_ocr(checks)
    check_obsidian_wiki(checks)
    check_command_surfaces(checks)
    check_docs_package(checks)

    payload, markdown = build_report(checks, report_id)
    json_path = out_dir / f"{report_id}.json"
    md_path = out_dir / f"{report_id}.md"
    last_json = out_dir / "pending-check-last.json"
    last_md = out_dir / "pending-check-last.md"
    json_text = json.dumps(payload, ensure_ascii=False, indent=2)
    json_path.write_text(json_text, encoding="utf-8")
    md_path.write_text(markdown, encoding="utf-8")
    last_json.write_text(json_text, encoding="utf-8")
    last_md.write_text(markdown, encoding="utf-8")

    print(markdown)
    print("")
    print(f"[REPORT_MD]   {md_path}")
    print(f"[REPORT_JSON] {json_path}")
    return 1 if payload["counts"]["FAIL"] else 0


if __name__ == "__main__":
    raise SystemExit(main())
