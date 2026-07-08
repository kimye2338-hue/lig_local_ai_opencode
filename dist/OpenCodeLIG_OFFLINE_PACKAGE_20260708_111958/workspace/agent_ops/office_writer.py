# -*- coding: utf-8 -*-
"""Office 파일 생성 (오프라인, Office 설치 불필요) — python-docx/pptx/openpyxl 래퍼.

우리 산출물은 지금까지 .md 초안뿐이고 실제 .docx/.pptx 는 "app validation pending"
(COM/앱 필요) 이었다. 순수 Python 라이브러리(openpyxl/python-docx/python-pptx)는
**Office 설치 없이** 진짜 파일을 오프라인으로 만든다 — wheel 만 반입하면 됨.

설계(우리 optional-dep 패턴 — doc_convert/ocr_screen 과 동일):
  - 라이브러리가 있으면 실제 파일 생성, 없으면 조용히 실패 않고 **반입 안내**.
  - 디자인 원칙(제목 위계·표 헤더 굵게·숫자 우측정렬·1슬라이드1메시지)을 반영.
  - COM 어댑터(excel_com/office_convert)와 별개: 이건 앱 없이 만드는 경로.

반입(오프라인): 인터넷 PC에서 `pip download python-docx python-pptx openpyxl -d wheelhouse`
→ 회사 PC에서 `pip install --no-index --find-links wheelhouse python-docx python-pptx openpyxl`.
"""
from __future__ import annotations

from pathlib import Path
from typing import Any, Dict, List, Optional, Sequence

_HINT = {
    "docx": "python-docx", "pptx": "python-pptx", "xlsx": "openpyxl",
}


def available(fmt: str) -> bool:
    mod = {"docx": "docx", "pptx": "pptx", "xlsx": "openpyxl"}.get(fmt)
    if not mod:
        return False
    try:
        __import__(mod)
        return True
    except Exception:
        return False


def _miss(fmt: str) -> Dict[str, Any]:
    lib = _HINT.get(fmt, fmt)
    return {"ok": False, "error": f"{lib} 미반입",
            "hint": f"오프라인 반입: pip download {lib} -d wheelhouse → "
                    f"pip install --no-index --find-links wheelhouse {lib} (docs/기능/OFFICE_WRITER.md)"}


def _is_number(s: Any) -> bool:
    try:
        float(str(s).replace(",", "").strip())
        return True
    except Exception:
        return False


def write_xlsx(path: Path, headers: Sequence[str], rows: Sequence[Sequence[Any]],
               sheet: str = "Sheet1") -> Dict[str, Any]:
    """표를 .xlsx 로. 헤더 굵게, 숫자열 우측정렬 + 숫자형 저장."""
    if not available("xlsx"):
        return _miss("xlsx")
    try:
        from openpyxl import Workbook  # type: ignore
        from openpyxl.styles import Font, Alignment  # type: ignore
        wb = Workbook()
        ws = wb.active
        ws.title = sheet[:31] or "Sheet1"
        if headers:
            ws.append(list(headers))
            for c in range(1, len(headers) + 1):
                ws.cell(row=1, column=c).font = Font(bold=True)
        for r in rows:
            out = []
            for cell in r:
                out.append(float(str(cell).replace(",", "")) if _is_number(cell) else cell)
            ws.append(out)
        # 숫자열 우측정렬
        for c in range(1, (len(headers) or (len(rows[0]) if rows else 0)) + 1):
            col_vals = [rows[i][c - 1] for i in range(len(rows)) if c - 1 < len(rows[i])]
            if col_vals and all(_is_number(v) for v in col_vals):
                for rr in range(2, len(rows) + 2):
                    ws.cell(row=rr, column=c).alignment = Alignment(horizontal="right")
        Path(path).parent.mkdir(parents=True, exist_ok=True)
        wb.save(str(path))
        return {"ok": True, "path": str(path), "format": "xlsx", "engine": "openpyxl"}
    except Exception as exc:  # noqa: BLE001
        return {"ok": False, "error": f"xlsx 생성 실패: {exc!r}", "hint": _miss('xlsx')['hint']}


def write_docx(path: Path, title: str, sections: Sequence[Dict[str, Any]]) -> Dict[str, Any]:
    """제목 + 섹션들을 .docx 로. section={heading, paragraphs[], bullets[], table:{headers,rows}}."""
    if not available("docx"):
        return _miss("docx")
    try:
        from docx import Document  # type: ignore
        doc = Document()
        doc.add_heading(title, level=0)
        for sec in sections:
            if sec.get("heading"):
                doc.add_heading(str(sec["heading"]), level=1)
            for p in sec.get("paragraphs", []) or []:
                doc.add_paragraph(str(p))
            for b in sec.get("bullets", []) or []:
                doc.add_paragraph(str(b), style="List Bullet")
            tbl = sec.get("table")
            if tbl and tbl.get("headers"):
                headers = tbl["headers"]
                rows = tbl.get("rows", [])
                t = doc.add_table(rows=1, cols=len(headers))
                try:
                    t.style = "Light Grid Accent 1"
                except Exception:
                    pass
                for i, h in enumerate(headers):
                    t.rows[0].cells[i].text = str(h)
                for r in rows:
                    cells = t.add_row().cells
                    for i, c in enumerate(r[:len(headers)]):
                        cells[i].text = str(c)
        Path(path).parent.mkdir(parents=True, exist_ok=True)
        doc.save(str(path))
        return {"ok": True, "path": str(path), "format": "docx", "engine": "python-docx"}
    except Exception as exc:  # noqa: BLE001
        return {"ok": False, "error": f"docx 생성 실패: {exc!r}", "hint": _miss('docx')['hint']}


def write_pptx(path: Path, slides: Sequence[Dict[str, Any]], title: str = "") -> Dict[str, Any]:
    """슬라이드들을 .pptx 로. slide={title(=핵심 메시지), points[]}. 1슬라이드=1메시지."""
    if not available("pptx"):
        return _miss("pptx")
    try:
        from pptx import Presentation  # type: ignore
        prs = Presentation()
        if title:
            s = prs.slides.add_slide(prs.slide_layouts[0])
            s.shapes.title.text = title
        bullet_layout = prs.slide_layouts[1]
        for sl in slides:
            s = prs.slides.add_slide(bullet_layout)
            s.shapes.title.text = str(sl.get("title", ""))[:120]
            body = s.placeholders[1].text_frame
            pts = sl.get("points", []) or []
            body.text = str(pts[0]) if pts else ""
            for p in pts[1:6]:  # 글머리 6줄 이내
                body.add_paragraph().text = str(p)
        Path(path).parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(path))
        return {"ok": True, "path": str(path), "format": "pptx", "engine": "python-pptx"}
    except Exception as exc:  # noqa: BLE001
        return {"ok": False, "error": f"pptx 생성 실패: {exc!r}", "hint": _miss('pptx')['hint']}


def status() -> Dict[str, bool]:
    return {fmt: available(fmt) for fmt in ("docx", "pptx", "xlsx")}


if __name__ == "__main__":
    import json
    print(json.dumps(status(), ensure_ascii=False, indent=2))
